import os


def ppt_to_pdf(ppt_path, output_path):
    """Convert a PowerPoint file (.pptx/.ppt) to PDF.
    Uses Windows COM automation (requires PowerPoint installed).
    Falls back to a slide-image approach if COM fails.

    NOTE: pythoncom/win32com are imported lazily INSIDE the try block
    below so this module can load on Linux/macOS without crashing the app.
    """
    try:
        # Windows COM approach - requires MS PowerPoint installed
        import pythoncom
        pythoncom.CoInitialize()
        import win32com.client
        
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = False
        
        abs_ppt = os.path.abspath(ppt_path)
        abs_out = os.path.abspath(output_path)
        
        presentation = powerpoint.Presentations.Open(abs_ppt, WithWindow=False)
        presentation.SaveAs(abs_out, 32)  # 32 = ppSaveAsPDF
        presentation.Close()
        powerpoint.Quit()
        
        pythoncom.CoUninitialize()
        return output_path
        
    except Exception:
        # Fallback: use python-pptx to extract text and create a simple PDF
        try:
            from pptx import Presentation
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import landscape, letter
            
            prs = Presentation(ppt_path)
            c = canvas.Canvas(output_path, pagesize=landscape(letter))
            width, height = landscape(letter)
            
            for slide_num, slide in enumerate(prs.slides):
                if slide_num > 0:
                    c.showPage()
                
                # Draw slide number
                c.setFont("Helvetica-Bold", 14)
                c.drawString(40, height - 50, f"Slide {slide_num + 1}")
                
                # Extract and render text from shapes
                y_pos = height - 80
                c.setFont("Helvetica", 11)
                
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                # Word wrap long lines
                                while len(text) > 90:
                                    c.drawString(40, y_pos, text[:90])
                                    text = text[90:]
                                    y_pos -= 16
                                    if y_pos < 50:
                                        c.showPage()
                                        y_pos = height - 50
                                
                                c.drawString(40, y_pos, text)
                                y_pos -= 16
                                
                                if y_pos < 50:
                                    c.showPage()
                                    y_pos = height - 50
            
            c.save()
            return output_path
            
        except Exception as e:
            raise RuntimeError(f"Could not convert PowerPoint to PDF: {e}")
