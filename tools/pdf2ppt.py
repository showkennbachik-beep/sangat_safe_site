import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import os

def pdf_to_pptx(pdf_path, output_path):
    doc = fitz.open(pdf_path)
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # fully blank layout
    out_dir = os.path.dirname(output_path)

    for page_num in range(len(doc)):
        page = doc[page_num]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        temp_img_path = os.path.join(out_dir, f"_tmp_ppt_{page_num}.png")
        pix.save(temp_img_path)

        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(
            temp_img_path, Inches(0), Inches(0),
            width=prs.slide_width, height=prs.slide_height
        )

        try:
            os.remove(temp_img_path)
        except OSError:
            pass

    doc.close()
    prs.save(output_path)
    return output_path
